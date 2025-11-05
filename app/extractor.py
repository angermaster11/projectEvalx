import base64
import io
from pptx import Presentation
from typing import List
from .schemas import SlideBlock

def _img_to_b64(image_blob: bytes) -> str:
    return "data:image/png;base64," + base64.b64encode(image_blob).decode("utf-8")

def extract_slides(ppt_bytes: bytes, max_slides: int = 60) -> List[SlideBlock]:
    prs = Presentation(io.BytesIO(ppt_bytes))
    slides = []
    for si, slide in enumerate(prs.slides):
        if si >= max_slides:
            break
        texts = []
        imgs = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                texts.append(shape.text)
            if shape.shape_type == 13 and hasattr(shape, "image"):
                try:
                    imgs.append(_img_to_b64(shape.image.blob))
                except:
                    pass
        slides.append(SlideBlock(index=si+1, text="\n".join([t.strip() for t in texts if t and t.strip()]), images=imgs))
    return slides
